"""제안 리스트(Excel) + 제안서(PPT) 생성 서비스.

템플릿 정의: .claude/skills/herofinder-proposal/SKILL.md
실제 대행사 제안 문서 구조를 일반화 — 브랜드/상품/인플루언서명은 전부 캠페인 데이터에서 주입.
개인 연락처는 산출물에 절대 포함하지 않는다.
"""

import uuid
from pathlib import Path

from loguru import logger
from openpyxl import Workbook
from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Inches, Pt

from ..models.orm_models import Campaign

OUTPUT_DIR = Path(__file__).parent.parent / "outputs"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

CHANNEL_SHEET_LABELS = {"youtube": "유튜브", "instagram": "인스타그램", "tiktok": "틱톡"}
AGENCY_NAME = "Hero Finder"
HIGHLIGHT_SCORE = 85  # 이 점수 이상 추천 후보 하이라이트

XLSX_HEADERS = [
    "No.", "카테고리", "채널명", "구독자수", "평균조회수", "URL", "진행 가능 여부",
    "진행 가능 일정", "비용", "비고", "브랜디드 진행 레퍼런스", "추천 사유", "예상 홍보 썸네일",
]

# 커뮤니케이션 퍼널 3단계 — 큐레이션 유형의 기본 축 (캠페인별 카테고리로 구체화)
FUNNEL_STAGES = [
    {"key": "인지", "role": "상품·브랜드를 나의 문제로 인식", "traits": ["현실 고민 공감", "쉬운 언어 설명", "미래 관점 제시"]},
    {"key": "행동", "role": "선택·행동 가이드 제공", "traits": ["시작 방법 제시", "실전 활용 팁", "행동 유도"]},
    {"key": "활용", "role": "심화 활용 전략 제시", "traits": ["비교 분석", "제도·트렌드 해석", "고급 활용 전략"]},
]

BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x73, 0x73, 0x73)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)


def _fmt_count(n: int) -> str:
    if n >= 100_000_000:
        return f"{n / 100_000_000:.1f}억"
    if n >= 10_000:
        return f"{n / 10_000:.1f}만"
    return str(n)


def _cost_manwon(rec: dict) -> str:
    lo, hi = rec.get("cost_range_min"), rec.get("cost_range_max")
    if not lo and not hi:
        return "협의"
    mid = ((lo or 0) + (hi or 0)) // 2 if (lo and hi) else (lo or hi or 0)
    return f"{mid // 10_000:,}"


def _reason_text(rec: dict) -> str:
    status = f"[채널 현황]\n{rec.get('category') or '-'} 카테고리 · 팔로워 {_fmt_count(rec['followers'])} · 참여율 {rec['engagement_rate']}% · 주간 성장률 +{rec['growth_rate']}%"
    proposal = f"[제안 사유]\n{rec.get('match_reason', '')}"
    return f"{status}\n\n{proposal}"


def generate_listup_xlsx(campaign: Campaign, recs: list[dict]) -> Path:
    """채널별 시트로 구성된 인플루언서 리스트업 엑셀 생성."""
    wb = Workbook()
    wb.remove(wb.active)

    thin = Border(*(Side(style="thin", color="D9D9D9"),) * 4)
    header_fill = PatternFill("solid", fgColor="1A1A1A")
    header_font = Font(color="FFFFFF", bold=True, size=10)
    green_fill = PatternFill("solid", fgColor="D9EAD3")  # 추천 후보 하이라이트
    wrap = Alignment(wrap_text=True, vertical="top")
    center = Alignment(horizontal="center", vertical="center", wrap_text=True)

    by_channel: dict[str, list[dict]] = {}
    for r in recs:
        by_channel.setdefault(r["channel"], []).append(r)

    for channel, rows in by_channel.items():
        ws = wb.create_sheet(f"<{CHANNEL_SHEET_LABELS.get(channel, channel)}>")
        ws.cell(row=1, column=11, value="*심의 관련 비용 추후 논의 필요").font = Font(size=9, color="808080")
        ws.cell(row=2, column=5, value="*최근 주간 스냅샷 기준").font = Font(size=9, color="808080")
        ws.cell(row=2, column=13, value=f"초록색: {AGENCY_NAME} 추천 후보").font = Font(size=9, color="38761D")

        for col, header in enumerate(XLSX_HEADERS, start=2):
            c = ws.cell(row=3, column=col, value=header)
            c.fill, c.font, c.alignment, c.border = header_fill, header_font, center, thin

        # 카테고리 기준 정렬 → 그룹 병합
        rows.sort(key=lambda r: (r.get("category") or "", -r["match_score"]))
        row_idx = 4
        cat_start: dict[str, int] = {}
        for i, rec in enumerate(rows, start=1):
            cat = rec.get("category") or "-"
            values = [
                i, cat, rec["name"], _fmt_count(rec["followers"]),
                _fmt_count(rec["monthly_views"]), rec.get("url") or "-",
                "O", "협의 가능", _cost_manwon(rec) + "만원",
                "*2차 라이선스 비용 별도 협의 필요",
                rec.get("reference_url") or "-", _reason_text(rec), "-",
            ]
            for col, v in enumerate(values, start=2):
                c = ws.cell(row=row_idx, column=col, value=v)
                c.border = thin
                c.alignment = wrap if col in (11, 13) else center
                if rec["match_score"] >= HIGHLIGHT_SCORE:
                    c.fill = green_fill
            cat_start.setdefault(cat, row_idx)
            row_idx += 1

        # 카테고리 세로 병합
        for cat, start in cat_start.items():
            end = start + sum(1 for r in rows if (r.get("category") or "-") == cat) - 1
            if end > start:
                ws.merge_cells(start_row=start, start_column=3, end_row=end, end_column=3)

        widths = {2: 5, 3: 12, 4: 16, 5: 10, 6: 10, 7: 34, 8: 12, 9: 12, 10: 10, 11: 26, 12: 30, 13: 44, 14: 14}
        for col, w in widths.items():
            ws.column_dimensions[ws.cell(row=3, column=col).column_letter].width = w

    out_path = OUTPUT_DIR / f"listup_{campaign.id}_{uuid.uuid4().hex[:8]}.xlsx"
    wb.save(out_path)
    return out_path


def _add_textbox(slide, left, top, width, height, text, size=18, bold=False, color=BLACK, align=None):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        if align:
            p.alignment = align
    return box


def _curation_types(recs: list[dict]) -> list[dict]:
    """추천 결과의 상위 카테고리를 퍼널 3단계와 결합해 큐레이션 유형 도출."""
    cat_counts: dict[str, int] = {}
    for r in recs:
        cat = r.get("category") or "기타"
        cat_counts[cat] = cat_counts.get(cat, 0) + 1
    top_cats = sorted(cat_counts, key=cat_counts.get, reverse=True)[:3]
    types = []
    for i, stage in enumerate(FUNNEL_STAGES):
        cat = top_cats[i] if i < len(top_cats) else (top_cats[-1] if top_cats else "일반")
        members = [r for r in recs if (r.get("category") or "기타") == cat][:4]
        types.append({
            "name": f"{cat} · {stage['key']} 중심형",
            "stage": stage,
            "members": members,
        })
    return types


def generate_proposal_pptx(campaign: Campaign, recs: list[dict]) -> Path:
    """클라이언트 내부 보고용 인플루언서 캠페인 기획안 PPT 생성 (16:9)."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    year = campaign.period_start.year if campaign.period_start else 2026

    def new_slide():
        return prs.slides.add_slide(blank)

    def fill_bg(slide, color):
        bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = color
        bg.line.fill.background()
        slide.shapes._spTree.remove(bg._element)
        slide.shapes._spTree.insert(2, bg._element)

    # 1. 표지 (블랙)
    s = new_slide()
    fill_bg(s, BLACK)
    _add_textbox(s, 1, 2.4, 11, 0.6, str(year), 20, color=GRAY)
    _add_textbox(s, 1, 3.0, 11, 1.0, campaign.name, 40, bold=True, color=WHITE)
    _add_textbox(s, 1, 4.1, 11, 0.6, "인플루언서 캠페인 기획안", 22, color=WHITE)
    _add_textbox(s, 1, 6.6, 11, 0.4, AGENCY_NAME, 14, color=GRAY)

    # 2. 운영 개요
    s = new_slide()
    _add_textbox(s, 0.8, 0.6, 11, 0.6, "운영 개요", 28, bold=True)
    overview = campaign.content_detail or "타깃별 관심사에 맞는 인플루언서 협업을 통해 인지도 및 관심도를 제고합니다."
    expectation = campaign.expectation or ""
    fmt = {"shortform": "숏폼", "longform": "롱폼", "package": "롱폼+숏폼 패키지"}.get(campaign.content_format or "", "협의")
    length = ""
    if campaign.longform_minutes:
        length += f" 롱폼 {campaign.longform_minutes}분"
    if campaign.shortform_minutes:
        length += f" 숏폼 {campaign.shortform_minutes}분"
    _add_textbox(s, 0.8, 1.7, 11.5, 1.2, overview, 16, color=GRAY)
    meta_lines = [
        f"광고 타입: {'PPL' if campaign.ad_type == 'ppl' else '브랜디드 콘텐츠'}   |   콘텐츠 규격: {fmt}{length}",
        f"예산: {campaign.budget_range or '협의'}   |   납품 기한: {campaign.deadline or '협의'}",
    ]
    if campaign.additional_rewards:
        meta_lines.append(f"추가 보상: {campaign.additional_rewards[:80]}")
    if campaign.must_include:
        meta_lines.append(f"필수 포함: {campaign.must_include[:80]}")
    if expectation:
        meta_lines.append(f"기대 효과: {expectation[:80]}")
    _add_textbox(s, 0.8, 3.2, 11.5, 2.5, "\n".join(meta_lines), 14)

    # 3. 채널 큐레이션 제안 (전략 요약)
    types = _curation_types(recs)
    s = new_slide()
    _add_textbox(s, 0.8, 0.6, 11.5, 0.6, "채널 큐레이션 제안", 28, bold=True)
    _add_textbox(s, 0.8, 1.4, 11.5, 0.5,
                 "채널별 특성에 따른 '인지 → 행동 → 활용' 커뮤니케이션으로 캠페인 목표 달성", 14, color=GRAY)
    for i, t in enumerate(types):
        left = 0.8 + i * 4.1
        card = s.shapes.add_shape(1, Inches(left), Inches(2.2), Inches(3.8), Inches(4.2))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_GRAY
        card.line.fill.background()
        _add_textbox(s, left + 0.25, 2.45, 3.3, 0.5, str(i + 1), 22, bold=True)
        _add_textbox(s, left + 0.25, 3.05, 3.3, 0.6, t["name"], 16, bold=True)
        examples = "\n".join(f"· {m['name']}" for m in t["members"][:3]) or "· (추천 채널)"
        _add_textbox(s, left + 0.25, 3.8, 3.3, 1.4, f"대표 추천 채널\n{examples}", 12, color=GRAY)
        _add_textbox(s, left + 0.25, 5.4, 3.3, 0.8, f"역할  |  {t['stage']['role']}", 12)

    # 4~6. 유형별 상세
    for i, t in enumerate(types):
        s = new_slide()
        _add_textbox(s, 0.8, 0.6, 11.5, 0.6, f"채널 큐레이션  |  {t['name']}", 24, bold=True)
        _add_textbox(s, 0.8, 1.5, 11.5, 0.5, f"{t['stage']['key']} 단계 — {t['stage']['role']}", 14, color=GRAY)
        _add_textbox(s, 0.8, 2.3, 3.0, 0.5, "채널 특징", 16, bold=True)
        for j, trait in enumerate(t["stage"]["traits"]):
            left = 0.8 + j * 4.1
            _add_textbox(s, left, 2.9, 3.8, 0.5, trait, 14, bold=True)
        _add_textbox(s, 0.8, 4.0, 3.0, 0.5, "추천 채널 & 제안 사유", 16, bold=True)
        y = 4.6
        for m in t["members"][:3]:
            _add_textbox(s, 0.8, y, 11.7, 0.8,
                         f"{m['name']}  ·  팔로워 {_fmt_count(m['followers'])}  ·  매칭 {m['match_score']}점\n{m.get('match_reason', '')[:110]}",
                         11, color=GRAY)
            y += 0.9

    # 7~8. 큐레이션 구성 리스트 테이블 (슬라이드당 최대 8행)
    table_headers = ["카테고리", "매체명", "구독자", "평균 조회수", "URL", "채널 현황", "추천 사유", "단가(만원)"]
    sorted_recs = sorted(recs, key=lambda r: -r["match_score"])
    for chunk_start in range(0, len(sorted_recs), 8):
        chunk = sorted_recs[chunk_start:chunk_start + 8]
        s = new_slide()
        _add_textbox(s, 0.8, 0.5, 11.5, 0.6, "채널 큐레이션 구성 리스트", 24, bold=True)
        tbl_shape = s.shapes.add_table(len(chunk) + 1, 8, Inches(0.6), Inches(1.4), Inches(12.1), Inches(0.5 + 0.62 * len(chunk)))
        tbl = tbl_shape.table
        col_widths = [1.2, 1.6, 1.0, 1.1, 2.2, 1.6, 2.6, 0.8]
        for c, w in enumerate(col_widths):
            tbl.columns[c].width = Inches(w)
        for c, h in enumerate(table_headers):
            cell = tbl.cell(0, c)
            cell.text = h
            cell.text_frame.paragraphs[0].font.size = Pt(11)
            cell.text_frame.paragraphs[0].font.bold = True
        for r, rec in enumerate(chunk, start=1):
            vals = [
                rec.get("category") or "-", rec["name"], _fmt_count(rec["followers"]),
                _fmt_count(rec["monthly_views"]), (rec.get("url") or "-")[:40],
                f"참여율 {rec['engagement_rate']}%", rec.get("match_reason", "")[:60],
                _cost_manwon(rec),
            ]
            for c, v in enumerate(vals):
                cell = tbl.cell(r, c)
                cell.text = str(v)
                cell.text_frame.paragraphs[0].font.size = Pt(9)
        _add_textbox(s, 0.6, 6.9, 12, 0.4,
                     f"* 하이라이트: {AGENCY_NAME} 추천 계정, 세부 사항은 리스트업 엑셀 참고 (2차 라이선스 별도 협의 필요)",
                     10, color=GRAY)

    # 9. E.O.D (블랙)
    s = new_slide()
    fill_bg(s, BLACK)
    _add_textbox(s, 1, 3.3, 11.3, 0.8, "E.O.D", 36, bold=True, color=WHITE)

    out_path = OUTPUT_DIR / f"proposal_{campaign.id}_{uuid.uuid4().hex[:8]}.pptx"
    prs.save(out_path)
    return out_path


def generate_proposal_package(campaign: Campaign, recs: list[dict]) -> dict:
    """리스트업 xlsx + 제안서 pptx를 함께 생성. 실패 시 개별 폴백."""
    result: dict = {"xlsx_path": None, "pptx_path": None}
    try:
        result["xlsx_path"] = f"/outputs/{generate_listup_xlsx(campaign, recs).name}"
    except Exception as e:
        logger.error(f"리스트업 엑셀 생성 실패 (campaign={campaign.id}): {e}")
    try:
        result["pptx_path"] = f"/outputs/{generate_proposal_pptx(campaign, recs).name}"
    except Exception as e:
        logger.error(f"제안서 PPT 생성 실패 (campaign={campaign.id}): {e}")
    return result
